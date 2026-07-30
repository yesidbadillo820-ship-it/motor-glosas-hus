"""
nucleo.office_tools — Conversión Office ⇄ PDF (fase 2 de la caja de herramientas).

Dos motores, según la dirección:

  Office → PDF  (Word/Excel/PowerPoint/HTML → PDF y PDF/A)
      Usa LibreOffice en modo headless (`soffice --convert-to`). Es el único
      camino fiable y por lotes; el analista instala LibreOffice una vez
      (gratis). Si no está, se avisa con instrucciones claras.

  PDF → Office  (PDF → Word / Excel / PowerPoint)
      Usa librerías Python que corren OFFLINE y NO necesitan LibreOffice:
        · PDF→Word   con pdf2docx   (respeta el diseño)
        · PDF→Excel  con pdfplumber (extrae tablas a hojas)
        · PDF→PPT    con PyMuPDF + python-pptx (una diapositiva por página)

  PDF → PDF/A  (archivado oficial) usa LibreOffice (importa el PDF en Draw).

Cada función devuelve la ruta de salida. `log` recibe el avance.
"""

import io
import os
import shutil
import subprocess
import tempfile

MENSAJE_SIN_LO = (
    "Para convertir documentos de Office a PDF hace falta LibreOffice "
    "(gratis). Instálelo desde https://es.libreoffice.org/descarga/ y vuelva a "
    "intentar. Las conversiones de PDF → Word/Excel/PowerPoint sí funcionan sin él."
)

# Filtro de exportación PDF según el tipo de origen (necesario para PDF/A).
_EXT_FILTRO = {
    ".pdf": "draw_pdf_Export",
    ".xlsx": "calc_pdf_Export",
    ".xls": "calc_pdf_Export",
    ".ods": "calc_pdf_Export",
    ".csv": "calc_pdf_Export",
    ".pptx": "impress_pdf_Export",
    ".ppt": "impress_pdf_Export",
    ".odp": "impress_pdf_Export",
}

# Extensiones de Office/HTML que LibreOffice sabe pasar a PDF.
EXT_OFFICE = (
    ".doc",
    ".docx",
    ".odt",
    ".rtf",
    ".txt",
    ".html",
    ".htm",
    ".xls",
    ".xlsx",
    ".ods",
    ".csv",
    ".ppt",
    ".pptx",
    ".odp",
)


# ------------------------------------------------------------- LibreOffice


def buscar_soffice():
    """Ubica el ejecutable de LibreOffice (PATH o rutas típicas de Windows)."""
    for nombre in ("soffice", "libreoffice"):
        ruta = shutil.which(nombre)
        if ruta:
            return ruta
    for p in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ):
        if os.path.exists(p):
            return p
    return None


def hay_libreoffice():
    return buscar_soffice() is not None


def _filtro_pdf(ext):
    return _EXT_FILTRO.get(ext.lower(), "writer_pdf_Export")


def a_pdf(ruta, carpeta=None, salida=None, pdfa=False, log=print):
    """Convierte Word/Excel/PowerPoint/HTML (o un PDF, para PDF/A) a PDF.

    pdfa=True produce un PDF/A-1b (archivable). Requiere LibreOffice.
    LibreOffice siempre escribe `base.pdf`; para no pisar el original (p. ej.
    un PDF→PDF/A en la misma carpeta) se convierte en un temporal y luego se
    mueve al destino final, que para PDF/A sobre un PDF lleva sufijo `_pdfa`.
    """
    soffice = buscar_soffice()
    if not soffice:
        raise RuntimeError(MENSAJE_SIN_LO)

    ext = os.path.splitext(ruta)[1]
    base = os.path.splitext(os.path.basename(ruta))[0]
    if salida:
        final = salida
    else:
        carpeta = carpeta or os.path.dirname(os.path.abspath(ruta))
        sufijo = "_pdfa" if (pdfa and ext.lower() == ".pdf") else ""
        final = os.path.join(carpeta, base + sufijo + ".pdf")
    os.makedirs(os.path.dirname(os.path.abspath(final)), exist_ok=True)

    destino = "pdf"
    if pdfa:
        destino = 'pdf:%s:{"SelectPdfVersion":{"type":"long","value":"1"}}' % _filtro_pdf(ext)

    tmp = tempfile.mkdtemp(prefix="lo_suite_")
    try:
        entorno = dict(os.environ)
        entorno.setdefault("HOME", tmp)  # LibreOffice necesita un HOME escribible
        cmd = [
            soffice,
            "-env:UserInstallation=file://%s" % tmp.replace(os.sep, "/"),
            "--headless",
            "--norestore",
            "--nologo",
            "--convert-to",
            destino,
            "--outdir",
            tmp,
            os.path.abspath(ruta),
        ]
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=180, env=entorno)
        generado = os.path.join(tmp, base + ".pdf")
        if not os.path.exists(generado):
            detalle = (proc.stderr or proc.stdout or "").strip().splitlines()
            pista = detalle[-1] if detalle else "sin detalle"
            raise RuntimeError(
                "LibreOffice no pudo convertir '%s': %s" % (os.path.basename(ruta), pista)
            )
        shutil.move(generado, final)
    finally:
        shutil.rmtree(tmp, ignore_errors=True)

    log("%s → %s%s" % (os.path.basename(ruta), final, "  (PDF/A)" if pdfa else ""))
    return final


def pdf_a_pdfa(ruta, carpeta=None, salida=None, log=print):
    """PDF normal → PDF/A-1b (para radicaciones que exigen archivable)."""
    return a_pdf(ruta, carpeta=carpeta, salida=salida, pdfa=True, log=log)


# ------------------------------------------------------------- PDF → Office
# (offline, sin LibreOffice)


def _reemplazar_ext(ruta, nueva_ext, carpeta=None):
    base = os.path.splitext(os.path.basename(ruta))[0]
    carpeta = carpeta or os.path.dirname(os.path.abspath(ruta))
    os.makedirs(carpeta, exist_ok=True)
    return os.path.join(carpeta, base + nueva_ext)


def pdf_a_word(ruta, salida=None, log=print):
    """PDF → Word (.docx) respetando el diseño, con pdf2docx."""
    from pdf2docx import Converter

    salida = salida or _reemplazar_ext(ruta, ".docx")
    cv = Converter(ruta)
    try:
        cv.convert(salida)
    finally:
        cv.close()
    log("PDF → Word: %s" % salida)
    return salida


def pdf_a_excel(ruta, salida=None, log=print):
    """PDF → Excel (.xlsx): una hoja por página, con sus tablas (o el texto)."""
    import pdfplumber
    from openpyxl import Workbook

    salida = salida or _reemplazar_ext(ruta, ".xlsx")
    wb = Workbook()
    wb.remove(wb.active)
    with pdfplumber.open(ruta) as pdf:
        for i, pagina in enumerate(pdf.pages, 1):
            ws = wb.create_sheet("Pagina_%d" % i)
            tablas = pagina.extract_tables()
            if tablas:
                for tabla in tablas:
                    for fila in tabla:
                        ws.append([celda if celda is not None else "" for celda in fila])
                    ws.append([])
            else:
                for linea in (pagina.extract_text() or "").splitlines():
                    ws.append([linea])
    if not wb.sheetnames:
        wb.create_sheet("Vacio")
    wb.save(salida)
    log("PDF → Excel: %s (%d hoja/s)" % (salida, len(wb.sheetnames)))
    return salida


def pdf_a_powerpoint(ruta, salida=None, dpi=120, log=print):
    """PDF → PowerPoint (.pptx): cada página como una diapositiva (imagen)."""
    import fitz
    from pptx import Presentation

    salida = salida or _reemplazar_ext(ruta, ".pptx")
    prs = Presentation()
    ancho = prs.slide_width
    doc = fitz.open(ruta)
    for pagina in doc:
        pix = pagina.get_pixmap(dpi=dpi)
        flujo = io.BytesIO(pix.tobytes("png"))
        diap = prs.slides.add_slide(prs.slide_layouts[6])  # en blanco
        alto = int(ancho * pix.height / pix.width)
        diap.shapes.add_picture(flujo, 0, 0, width=ancho, height=alto)
    doc.close()
    prs.save(salida)
    log("PDF → PowerPoint: %s (%d diapositiva/s)" % (salida, len(prs.slides._sldIdLst)))
    return salida
